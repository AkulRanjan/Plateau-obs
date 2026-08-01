"""Build the jury deck.

    python demo/make_deck.py            # -> demo/Plateau-jury-deck.pptx

Every number in here is measured and traceable to something in this repo — the
two-laptop run's JSONL logs, metrics.json, or the test suite. Nothing is
projected, illustrative or rounded in our favour, except the token figures,
which are arithmetic over an assumed 1,850/turn and are labelled as such on the
slide itself.

If a number changes, change it HERE and regenerate. Do not edit the .pptx by
hand — a deck that disagrees with the repo is the one thing a judge who opens
the repo will find.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).resolve().parent / "Plateau-jury-deck.pptx"

# The product's palette: chart stock, plotter ink, one signal red.
PAPER = RGBColor(0xE6, 0xE3, 0xDA)
DESK = RGBColor(0xDC, 0xD7, 0xC9)
INK = RGBColor(0x1C, 0x1B, 0x18)
MUTED = RGBColor(0x4F, 0x4B, 0x42)
FAINT = RGBColor(0x6E, 0x69, 0x5C)
PEN_A = RGBColor(0x1B, 0x4B, 0x8F)
PEN_B = RGBColor(0x0E, 0x6F, 0x52)
SIGNAL = RGBColor(0xB3, 0x22, 0x1B)
AMBER = RGBColor(0x8A, 0x5F, 0x0B)

# Arial Narrow and Consolas are the two faces present on essentially every
# machine that will open this, and LibreOffice substitutes both sensibly.
DISPLAY = "Arial Narrow"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
L = Inches(0.72)          # left margin
CW = W - L * 2            # content width


def deck() -> Presentation:
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def slide(prs, eyebrow: str | None = None, title: str | None = None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = PAPER
    y = Inches(0.52)
    if eyebrow:
        text(s, L, y, CW, Inches(0.26), eyebrow, DISPLAY, 12, FAINT, space=0.18, caps=True)
        y += Inches(0.30)
    if title:
        text(s, L, y, CW, Inches(0.7), title, DISPLAY, 34, INK, bold=True, space=0.02)
        y += Inches(0.74)
        rule(s, L, y, CW, INK, Pt(1.5))
        y += Inches(0.18)
    return s, y


def rule(s, x, y, w, colour, thickness=Pt(0.75)):
    from pptx.enum.shapes import MSO_SHAPE

    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, thickness)
    ln.fill.solid()
    ln.fill.fore_color.rgb = colour
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def box(s, x, y, w, h, fill=DESK, edge=INK):
    from pptx.enum.shapes import MSO_SHAPE

    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    b.fill.solid()
    b.fill.fore_color.rgb = fill
    b.line.color.rgb = edge
    b.line.width = Pt(0.75)
    b.shadow.inherit = False
    return b


def text(s, x, y, w, h, body, font=DISPLAY, size=16, colour=INK,
         bold=False, space=0.0, caps=False, align=PP_ALIGN.LEFT, line=1.25):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(str(body).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        r = p.add_run()
        r.text = para.upper() if caps else para
        f = r.font
        f.name, f.size, f.bold = font, Pt(size), bold
        f.color.rgb = colour
        if space:
            _spacing(r, space)
    return tb


def _spacing(run, em: float):
    """Letter-spacing, which python-pptx does not expose."""
    run.font._rPr.set("spc", str(int(em * run.font.size.pt * 100)))


def stat(s, x, y, w, label, value, sub, value_colour=INK, h=Inches(1.28)):
    box(s, x, y, w, h)
    text(s, x + Inches(0.16), y + Inches(0.14), w - Inches(0.3), Inches(0.2),
         label, DISPLAY, 10.5, FAINT, space=0.14, caps=True)
    text(s, x + Inches(0.16), y + Inches(0.40), w - Inches(0.3), Inches(0.5),
         value, MONO, 30, value_colour, bold=True)
    text(s, x + Inches(0.16), y + Inches(0.95), w - Inches(0.3), Inches(0.28),
         sub, DISPLAY, 11, MUTED)


def table(s, x, y, w, rows, widths, head_fill=INK, size=12.5):
    """A plain ruled table. python-pptx's styling is fought rather than used."""
    n_rows, n_cols = len(rows), len(rows[0])
    row_h = Inches(0.42)
    shape = s.shapes.add_table(n_rows, n_cols, x, y, w, row_h * n_rows)
    t = shape.table
    t.first_row = True
    for i, cw in enumerate(widths):
        t.columns[i].width = Emu(int(w * cw))
    for r, row in enumerate(rows):
        t.rows[r].height = row_h
        for c, cell in enumerate(row):
            val, colour, bold = (cell if isinstance(cell, tuple) else (cell, INK, False))
            tc = t.cell(r, c)
            tc.fill.solid()
            tc.fill.fore_color.rgb = head_fill if r == 0 else (DESK if r % 2 else PAPER)
            tc.margin_left = tc.margin_right = Inches(0.12)
            tc.margin_top = tc.margin_bottom = Inches(0.05)
            tf = tc.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if c and r else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            f = run.font
            f.name = MONO if (c and r) else DISPLAY
            f.size = Pt(size)
            f.bold = bold or r == 0
            f.color.rgb = PAPER if r == 0 else colour
    return shape


def note(s, y, body):
    text(s, L, y, CW, Inches(0.6), body, DISPLAY, 12, FAINT, line=1.35)


# ---------------------------------------------------------------------------


def build() -> Presentation:
    prs = deck()

    # 1 -------------------------------------------------------------- title
    s, _ = slide(prs)
    rule(s, L, Inches(2.5), Inches(1.6), SIGNAL, Pt(3))
    text(s, L, Inches(2.62), CW, Inches(1.3), "PLATEAU", DISPLAY, 76, INK,
         bold=True, space=0.22)
    text(s, L, Inches(4.28), Inches(9.4), Inches(0.9),
         "An agent that gets stuck does not crash. It bills.",
         DISPLAY, 27, MUTED)
    rule(s, L, Inches(5.3), CW, INK, Pt(1))
    text(s, L, Inches(5.5), Inches(7.5), Inches(0.9),
         "A semantic circuit breaker for autonomous agents\n"
         "Team ABCD  ·  AI Safety and Observability",
         DISPLAY, 14, MUTED, line=1.5)
    text(s, L + Inches(8.6), Inches(5.5), Inches(3.3), Inches(0.9),
         "Measured live on two laptops\n137 tests  ·  runs fully offline",
         MONO, 12, FAINT, line=1.5)

    # 2 ------------------------------------------------------------ problem
    s, y = slide(prs, "The problem", "Nothing you already run will catch this.")
    text(s, L, y + Inches(0.1), Inches(6.4), Inches(2.4),
         "An agent loops on the same dead end. It rephrases the request, calls the "
         "same tool, gets the same non-answer, and tries again.\n\n"
         "Every call returns 200 OK. Error rate zero. Latency normal. Every "
         "dashboard green.", DISPLAY, 17, INK, line=1.45)
    b = box(s, L + Inches(7.0), y + Inches(0.1), Inches(4.9), Inches(2.5), fill=DESK)
    text(s, L + Inches(7.2), y + Inches(0.3), Inches(4.5), Inches(2.0),
         "monitoring today answers\n"
         "   did the call succeed?\n\n"
         "nobody answers\n"
         "   is the sequence going anywhere?",
         MONO, 14, MUTED, line=1.6)
    rule(s, L, y + Inches(2.95), CW, INK, Pt(1))
    text(s, L, y + Inches(3.15), CW, Inches(0.6),
         "The only thing that eventually notices is the invoice.",
         DISPLAY, 22, SIGNAL, bold=True)

    # 3 ------------------------------------------------------------- result
    s, y = slide(prs, "The result — measured, not projected",
                 "Two laptops. Two real agents. One impossible task.")
    rows = [
        ["", "Unguarded", "With Plateau"],
        ["Turns run", "25", "21"],
        ["Turns refused before reaching a tool", "0", ("4", PEN_B, True)],
        ["Turns spent re-asking a dead question",
         ("18", SIGNAL, True), ("6", PEN_B, True)],
        ["Tokens (est., 1,850/turn)", "46,250", "38,850"],
        ["Stopped by", ("we stopped it", SIGNAL, False), ("itself, turn 12", PEN_B, True)],
    ]
    table(s, L, y + Inches(0.05), CW, rows, [0.46, 0.27, 0.27])
    note(s, y + Inches(2.85),
         "Same model, same seed, same tools, same task. The only difference is that the "
         "right-hand agent had a breaker attached.\nThe unguarded agent did not finish — it "
         "hit a turn cap we imposed. Left alone it does not stop.")
    rule(s, L, y + Inches(3.75), CW, INK, Pt(1))
    text(s, L, y + Inches(3.95), CW, Inches(0.6),
         "67% fewer turns spent going nowhere.", DISPLAY, 26, INK, bold=True)

    # 4 --------------------------------------------------------- the metric
    s, y = slide(prs, "Which number is honest",
                 "We lead on wasted turns, not tokens.")
    stat(s, L, y + Inches(0.1), Inches(3.7), "Tokens saved", "16%",
         "diluted — and capped by us", MUTED)
    stat(s, L + Inches(4.0), y + Inches(0.1), Inches(3.7), "Wasted turns cut", "67%",
         "survives the turn cap", PEN_B)
    stat(s, L + Inches(8.0), y + Inches(0.1), Inches(3.9), "Guard overhead", "0",
         "embeddings on refused turns", PEN_B)
    text(s, L, y + Inches(1.75), CW, Inches(2.2),
         "The token gap understates the result for two reasons, and we would rather say so "
         "than quote the bigger-sounding number.\n\n"
         "One — the guarded agent spent its freed turns doing real work, which is the point, "
         "not a cost.\n"
         "Two — the unguarded agent's total is bounded only by a cap we imposed. Its real "
         "number is unbounded.",
         DISPLAY, 16, INK, line=1.5)

    # 5 ------------------------------------------------------- not just off
    s, y = slide(prs, "It does not just stop the agent",
                 "It says why, and where to go instead.")
    box(s, L, y + Inches(0.1), CW, Inches(2.5), fill=DESK, edge=SIGNAL)
    text(s, L + Inches(0.22), y + Inches(0.28), CW - Inches(0.5), Inches(0.3),
         "Breaker open", DISPLAY, 12, SIGNAL, bold=True, space=0.18, caps=True)
    text(s, L + Inches(0.22), y + Inches(0.62), CW - Inches(0.5), Inches(0.7),
         "The same question is being asked repeatedly and the answers are not new. "
         "Change the source or the approach, not the wording.\n"
         "action_sim 1.000 >= ceiling 0.761;  obs_novelty 0.000 < floor 0.30",
         MONO, 12.5, INK, line=1.5)
    text(s, L + Inches(0.22), y + Inches(1.52), CW - Inches(0.5), Inches(0.5),
         "Escape: turn 15 produced the most new information in this run (novelty 0.896). "
         "Return to that result and build on it instead of re-asking.",
         MONO, 12.5, AMBER, line=1.5)
    text(s, L, y + Inches(2.85), CW, Inches(1.2),
         "A step cap can only say \"stop\". This is the difference between an alert and a "
         "diagnosis — the agent is handed a reason and a next move, in its own context, "
         "and can act on both.", DISPLAY, 17, INK, line=1.45)

    # 6 -------------------------------------------------------- not a cap
    s, y = slide(prs, "Why not just cap the steps",
                 "A cap is wrong in both directions.")
    rows = [
        ["Approach", "Catches the stall", "False-trips healthy work"],
        ["Step cap at 25", "at turn 27 — too late", ("yes", SIGNAL, True)],
        ["Any-repeat detector", "at turn 5", ("yes, everything", SIGNAL, True)],
        ["Exact-match detector", ("never fires", SIGNAL, True), "no"],
        ["Plateau", ("at turn 12", PEN_B, True), ("no", PEN_B, True)],
    ]
    table(s, L, y + Inches(0.05), CW, rows, [0.34, 0.33, 0.33])
    note(s, y + Inches(2.45),
         "Measured over five 16-61 turn traces, in metrics.json. Healthy work means a batch job "
         "that legitimately repeats the same call on new data — an invoice run, a poller. "
         "Stopping those is worse than never having a detector.")
    rule(s, L, y + Inches(3.35), CW, INK, Pt(1))
    text(s, L, y + Inches(3.55), CW, Inches(0.7),
         "It also recovers. Tripped at 12, found new information at 14, resumed, and only\n"
         "re-opened when the agent stalled a second time.",
         DISPLAY, 17, INK, line=1.4)

    # 7 --------------------------------------------------------- integration
    s, y = slide(prs, "Adoption cost", "Two hooks around a tool call.")
    box(s, L, y + Inches(0.1), Inches(6.4), Inches(2.95), fill=DESK)
    text(s, L + Inches(0.25), y + Inches(0.32), Inches(6.0), Inches(2.5),
         "before_tool(action)\n"
         "    veto, once the breaker is open\n"
         "    costs nothing while closed\n\n"
         "after_tool(action, observation)\n"
         "    score the completed turn\n"
         "    this is the call that embeds",
         MONO, 14, INK, line=1.55)
    text(s, L + Inches(7.0), y + Inches(0.1), Inches(4.9), Inches(2.6),
         "That is the whole surface. It is the same shape as a LangChain callback pair or "
         "a Strands BeforeToolCall / AfterToolCall adapter.\n\n"
         "Working LangGraph example in the repo: the agent logic does not know it is being "
         "watched.", DISPLAY, 16, INK, line=1.5)
    rule(s, L, y + Inches(3.25), CW, INK, Pt(1))
    text(s, L, y + Inches(3.45), CW, Inches(0.6),
         "No retraining. No proxy. No vendor lock. Runs entirely offline.",
         DISPLAY, 21, INK, bold=True)

    # 8 ------------------------------------------------------------ honesty
    s, y = slide(prs, "What we measured that did not work",
                 "Three findings we are not hiding.")
    items = [
        ("It does not catch a reworded stall.",
         "Differently-worded non-answers score 0.44-1.01 novelty; none fall below the floor. "
         "It works here because real search APIs return a consistent no-results string."),
        ("A lexical baseline matches our recall.",
         "On current traces it matches recall with no false trips and detects sooner. Our "
         "defensible claim is the one above: we hold on healthy repetitive work."),
        ("One of our two dials earns its keep; the other has not yet.",
         "Ablating action-similarity changes detection on none of the five trace classes — "
         "only how fast. We report it rather than describe a two-dial design we cannot evidence."),
    ]
    yy = y + Inches(0.1)
    for head, body in items:
        rule(s, L, yy, Inches(0.5), SIGNAL, Pt(2.5))
        text(s, L, yy + Inches(0.16), CW, Inches(0.3), head, DISPLAY, 17, INK, bold=True)
        text(s, L, yy + Inches(0.52), CW - Inches(0.4), Inches(0.55), body,
             DISPLAY, 13.5, MUTED, line=1.4)
        yy += Inches(1.18)
    note(s, yy + Inches(0.05),
         "Every one of these is measured in the repo's own metrics.json. A demo whose numbers "
         "contradict its repository is the one thing a judge who opens it will find.")

    # 9 ---------------------------------------------------------------- ask
    s, y = slide(prs, "Where it stands", "Built, measured, and running.")
    for i, (label, value, sub) in enumerate([
        ("Tests passing", "137", "including the live demo stack"),
        ("Machines proven on", "2", "Fedora + macOS, over LAN"),
        ("External services", "0", "no cloud, no API key"),
        ("Integration surface", "2 hooks", "around any tool call"),
    ]):
        stat(s, L + Inches(3.06) * i, y + Inches(0.1), Inches(2.86), label, value, sub,
             PEN_B if i else INK)
    rule(s, L, y + Inches(1.75), CW, INK, Pt(1))
    text(s, L, y + Inches(1.98), Inches(7.0), Inches(1.6),
         "Next\n"
         "· a trace class where the second dial earns its place, or drop it\n"
         "· paraphrase-robust novelty, the one gap we have measured and named\n"
         "· adapters shipped for LangGraph and Strands",
         DISPLAY, 16, INK, line=1.6)
    box(s, L + Inches(7.7), y + Inches(1.98), Inches(4.2), Inches(1.95))
    text(s, L + Inches(7.9), y + Inches(2.2), Inches(3.8), Inches(1.2),
         "Watch it live\n\ndemo/run_demo.sh\n\nboth agents, one dashboard",
         MONO, 13, MUTED, line=1.6)

    return prs


if __name__ == "__main__":
    build().save(OUT)
    print(f"wrote {OUT}  ({OUT.stat().st_size // 1024} KB)")
